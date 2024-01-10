from pptx import Presentation
from docx import Document
from docx.shared import Inches

# extract image
def extract_images_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    image_list = []

    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                image = shape.image
                image_bytes = image.blob
                image_list.append(image_bytes)

    return image_list

# display image
def display_images_in_docx(images, docx_file):
    document = Document()

    for image_bytes in images:
        image_path = 'temp_image.png'
        with open(image_path, 'wb') as img_file:
            img_file.write(image_bytes)

        document.add_picture(image_path, width=Inches(4))
 
    document.save(docx_file)

if __name__ == "__main__":
    ppt_file = "Unit1.pptx"
    docx_file = "Captured.docx"

    images = extract_images_from_ppt(ppt_file)
    display_images_in_docx(images, docx_file)
